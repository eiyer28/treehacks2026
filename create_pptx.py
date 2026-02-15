from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
BG_DARK = RGBColor(0x0F, 0x0F, 0x0F)
BG_SECTION = RGBColor(0x1A, 0x1A, 0x2E)
NVIDIA_GREEN = RGBColor(0x76, 0xB9, 0x00)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_RED = RGBColor(0xFF, 0x4C, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
TABLE_HEADER_BG = RGBColor(0x1E, 0x1E, 0x3A)
TABLE_ROW_BG = RGBColor(0x14, 0x14, 0x24)
TABLE_ALT_BG = RGBColor(0x1A, 0x1A, 0x30)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    return table_shape.table


def style_cell(cell, text, font_size=14, color=WHITE, bold=False, bg_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    cell.text_frame.word_wrap = True
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_accent_line(slide, left, top, width, color=NVIDIA_GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_device_box(slide, left, top, width, height, title, subtitle, details, border_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    add_textbox(slide, left + 0.15, top + 0.1, width - 0.3, 0.4, title, font_size=16, color=border_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.15, top + 0.5, width - 0.3, 0.35, subtitle, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.15, top + 0.85, width - 0.3, height - 1.0, details, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 3, 1.8, 7.333)
add_textbox(slide, 1, 2.0, 11.333, 1.5, "Edge Rescue", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.2, 11.333, 0.8, "Humans, Robots, and AI \u2014 Working Together", font_size=28, color=NVIDIA_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.1, 11.333, 0.7,
    "Sim-validated autonomous robotics for disaster response.",
    font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.6, 11.333, 0.5,
    "Plan. Simulate. Validate. Build. \u2014 No cloud, no internet, fully on edge hardware.",
    font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.5, 11.333, 0.5, "Eashan Iyer  |  Samuel Lihn  |  Gordon Jin  |  Trevor Kwan", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 6.1, 11.333, 0.5, "TreeHacks 2026", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Story (condensed)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "HATAY, TURKEY  \u2014  FEBRUARY 6, 2023", font_size=14, color=ACCENT_RED, bold=True)
add_accent_line(slide, 0.8, 0.8, 2, color=ACCENT_RED)

# Left: the disaster facts
tf = add_textbox(slide, 0.8, 1.2, 6.0, 5.5, "", font_size=18)
add_paragraph(tf, "4:17 AM. A 7.8-magnitude earthquake hits.", font_size=24, color=WHITE, bold=True)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "59,000 dead across Turkey and Syria.", font_size=20, color=ACCENT_RED, bold=True, space_before=8)
add_paragraph(tf, "230,000+ buildings damaged or destroyed.", font_size=18, color=WHITE, space_before=8)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "The flaw wasn't bravery. It was information.", font_size=20, color=WHITE, bold=True, space_before=16)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "Rescue teams couldn't reach remote villages for 3\u20134 days.", font_size=16, color=LIGHT_GRAY, space_before=6)
add_paragraph(tf, "Syria's White Helmets had no heavy equipment \u2014 it never came.", font_size=16, color=LIGHT_GRAY, space_before=6)
add_paragraph(tf, "At collapsed buildings, teams faced an impossible question:", font_size=16, color=LIGHT_GRAY, space_before=6)
add_paragraph(tf, "\"If we move this slab, will the rest collapse on us?\"", font_size=18, color=WHITE, bold=True, space_before=6)
add_paragraph(tf, "No simulation. No structural model. Just gut instinct", font_size=16, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "under extreme time pressure.", font_size=16, color=LIGHT_GRAY, space_before=0)

# Right: the survival clock
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.2), Inches(5.0), Inches(4.8)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x0A)
shape.line.color.rgb = ACCENT_RED
shape.line.width = Pt(2)

add_textbox(slide, 7.7, 1.4, 4.6, 0.5, "THE SURVIVAL CLOCK", font_size=16, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

clock_data = [
    ("< 24 hours", "90%", "survival rate"),
    ("24\u201348 hours", "50\u201360%", "survival rate"),
    ("48\u201372 hours", "20\u201330%", "survival rate"),
    ("> 72 hours", "5\u201310%", "survival rate"),
]
cy = 2.1
for time, rate, label in clock_data:
    add_textbox(slide, 7.9, cy, 2.0, 0.45, time, font_size=15, color=LIGHT_GRAY, bold=True)
    add_textbox(slide, 9.9, cy, 1.5, 0.45, rate, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, 11.5, cy + 0.05, 0.8, 0.4, label, font_size=11, color=LIGHT_GRAY)
    cy += 0.55

add_textbox(slide, 7.7, 4.4, 4.6, 0.9,
    "Every hour a robot waits for a structural assessment is an hour a survivor may not have.",
    font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom transition
tf = add_textbox(slide, 0.8, 6.3, 11.7, 0.9,
    "What if the robot didn't have to guess? What if it could simulate the physics of a collapse,",
    font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "test 100 approaches in seconds, and only commit to the safest one \u2014 with no internet?",
    font_size=20, color=NVIDIA_GREEN, bold=True, alignment=PP_ALIGN.CENTER, space_before=2)

# ============================================================
# SLIDE 3: Our Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "EDGE RESCUE", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)
tf = add_textbox(slide, 0.8, 1.0, 11.7, 0.8,
    "Plan. Simulate. Validate. Build.",
    font_size=36, color=WHITE, bold=True)

steps = [
    ("1", "PROMPT", "\"Build a reinforced shelter\" \u2014 natural language, spoken or typed", ACCENT_BLUE),
    ("2", "PLAN", "Local LLM on DGX Spark decomposes into subtasks \u2014 no cloud, no internet", ACCENT_BLUE),
    ("3", "SIMULATE", "Isaac Sim tests each action: will it fall? Will it survive lateral forces?", NVIDIA_GREEN),
    ("4", "VALIDATE", "Only physics-validated actions are sent to the real robot", NVIDIA_GREEN),
    ("5", "BUILD", "SO-101 arm executes the safe plan", WHITE),
    ("6", "ADAPT", "If reality doesn't match the sim (block slips, misalignment), re-plan automatically", ACCENT_RED),
]

y = 2.2
for num, title, desc, color in steps:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.55), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf_shape = shape.text_frame
    tf_shape.paragraphs[0].text = num
    tf_shape.paragraphs[0].font.size = Pt(18)
    tf_shape.paragraphs[0].font.color.rgb = BG_DARK
    tf_shape.paragraphs[0].font.bold = True
    tf_shape.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.8, y - 0.05, 2.0, 0.55, title, font_size=18, color=color, bold=True)
    add_textbox(slide, 3.8, y, 8.7, 0.55, desc, font_size=16, color=LIGHT_GRAY)
    y += 0.75

add_textbox(slide, 0.8, 6.8, 11.7, 0.5,
    "The AI breaks it in simulation so it doesn't break in reality.",
    font_size=22, color=NVIDIA_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: You May Have Seen... (Differentiation)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "YOU MAY HAVE SEEN...", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)

# Left column header
add_textbox(slide, 0.8, 1.1, 5.8, 0.5, "What Exists", font_size=22, color=LIGHT_GRAY, bold=True)

# DeepMind
tf = add_textbox(slide, 0.8, 1.8, 5.8, 1.2, "", font_size=16)
add_paragraph(tf, "Google DeepMind \u2014 Gemini Robotics 1.5", font_size=16, color=WHITE, bold=True)
add_paragraph(tf, "100B+ parameter VLA that \"thinks before it acts.\"", font_size=14, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "Requires cloud, massive compute, and internet.", font_size=14, color=ACCENT_RED, space_before=2)
add_paragraph(tf, "Unusable in a disaster zone.", font_size=14, color=ACCENT_RED, bold=True, space_before=2)

# Digital Twins
tf = add_textbox(slide, 0.8, 3.5, 5.8, 1.0, "", font_size=16)
add_paragraph(tf, "Industry Digital Twins (Siemens, BMW, NVIDIA Omniverse)", font_size=16, color=WHITE, bold=True)
add_paragraph(tf, "Mirror a factory floor for monitoring. Passive \u2014", font_size=14, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "they watch, they don't plan or validate.", font_size=14, color=ACCENT_RED, space_before=2)

# Text-to-Robot
tf = add_textbox(slide, 0.8, 4.9, 5.8, 1.0, "", font_size=16)
add_paragraph(tf, "Text-to-Robot (VoxPoser, RoboGPT)", font_size=16, color=WHITE, bold=True)
add_paragraph(tf, "Prompt \u2192 one-shot execution. No physics validation.", font_size=14, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "If the LLM hallucinates an unstable structure, the robot builds it anyway.", font_size=14, color=ACCENT_RED, space_before=2)

# Divider line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.1), Pt(2), Inches(5.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
shape.line.fill.background()

# Right column header
add_textbox(slide, 7.2, 1.1, 5.8, 0.5, "What We Built", font_size=22, color=NVIDIA_GREEN, bold=True)

right_items = [
    ("Runs entirely on local edge hardware", "DGX Spark + Jetson + Isaac Sim. No internet required."),
    ("The sim doesn't just mirror \u2014 it plans ahead", "Proposes the next construction phase and stress-tests it before the robot commits."),
    ("Every action is physics-validated", "The LLM proposes, the sim disposes. Nothing executes without passing physics checks."),
    ("Bidirectional feedback loop", "When reality doesn't match sim (block slips), the system re-plans autonomously."),
]

ry = 1.8
for title, desc in right_items:
    add_textbox(slide, 7.4, ry, 5.4, 0.35, title, font_size=16, color=NVIDIA_GREEN, bold=True)
    add_textbox(slide, 7.4, ry + 0.35, 5.4, 0.45, desc, font_size=14, color=LIGHT_GRAY)
    ry += 1.1

# Bottom punchline
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.75)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
shape.line.color.rgb = NVIDIA_GREEN
shape.line.width = Pt(2)

add_textbox(slide, 1.7, 6.55, 9.9, 0.65,
    "DeepMind showed that robots should think before they act. We showed you can do it on edge hardware with no internet.",
    font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: The Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "THE DEMO \u2014 EARTHQUAKE TEST", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)

demo_steps = [
    ("\u2776  PROMPT", "\"Build a tower\" \u2014 DGX Spark LLM breaks this into pick-and-place subtasks", ACCENT_BLUE),
    ("\u2777  BUILD", "Physical arm stacks blocks. Isaac Sim mirrors every move in real time (split screen).", WHITE),
    ("\u2778  STRESS TEST", "Isaac Sim applies lateral forces (earthquake simulation). Tower collapses in sim. Real tower is untouched.", ACCENT_RED),
    ("\u2779  REDESIGN", "LLM analyzes the simulated collapse. Redesigns with wider base and interlocking pattern.", NVIDIA_GREEN),
    ("\u277A  REBUILD", "Arm rebuilds the reinforced version. Isaac Sim confirms it survives the same earthquake.", NVIDIA_GREEN),
]

y = 1.2
for title, desc, color in demo_steps:
    add_textbox(slide, 0.8, y, 3.5, 0.8, title, font_size=22, color=color, bold=True)
    add_textbox(slide, 4.3, y + 0.05, 8.2, 0.8, desc, font_size=16, color=LIGHT_GRAY)
    y += 1.05

tf = add_textbox(slide, 0.8, 6.6, 11.7, 0.7,
    "The same loop that reinforces a block tower is what would prevent a rescue robot from pulling the wrong beam.",
    font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "ARCHITECTURE", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)
add_textbox(slide, 0.8, 1.0, 11.7, 0.7,
    "Three devices. Four agents. Zero cloud dependency.",
    font_size=30, color=WHITE, bold=True)

# Top: Voice/Text prompt
add_textbox(slide, 4.5, 1.7, 4.3, 0.45, "\u25BC  Voice / Text Prompt  \u25BC", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# DGX Spark
add_device_box(slide, 4.2, 2.2, 4.9, 1.5,
    "DGX SPARK", "Agent C: The Overlord",
    "Local LLM \u2022 Task decomposition \u2022 Orchestration", NVIDIA_GREEN)

# Arrow down left
add_textbox(slide, 3.0, 3.75, 2.5, 0.4, "subtasks \u25BC", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
# Arrow down right
add_textbox(slide, 7.8, 3.75, 2.5, 0.4, "validation requests \u25BC", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Jetson (bottom left)
add_device_box(slide, 0.8, 4.2, 5.0, 1.8,
    "JETSON ORIN NANO", "Agent A: The Builder",
    "VLM perception \u2022 Motor control \u2022 Camera feed\n\u2192 Drives SO-101 Follower Arm", ACCENT_BLUE)

# Isaac Sim (bottom right)
add_device_box(slide, 7.5, 4.2, 5.0, 1.8,
    "RTX 3070 \u2014 ISAAC SIM", "Agent B: The Architect",
    "Physics simulation \u2022 Stress testing \u2022 Validation\n\u2192 Digital Twin of workspace", NVIDIA_GREEN)

# Feedback arrow
add_textbox(slide, 4.0, 6.1, 5.3, 0.4, "\u25C0\u2500\u2500 feedback loop \u2500\u2500\u25B6", font_size=14, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# Agent D note
add_textbox(slide, 0.8, 6.6, 11.7, 0.6,
    "Agent D: The Operator \u2014 Human override via Leader Arm for safety-critical fallback",
    font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: Why Edge + Hardware
# (Merged old "Why Edge Matters" and "Hardware" slides)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "WHY EDGE MATTERS", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)
tf = add_textbox(slide, 0.8, 1.0, 11.7, 0.6,
    "Hardware-native. Designed for environments where cloud doesn't exist.",
    font_size=26, color=WHITE, bold=True)

# Comparison table (compact)
tbl = add_table(slide, 6, 3, 0.8, 1.8, 6.0, 3.5)
headers = ["", "Cloud", "Edge Rescue"]
for i, h in enumerate(headers):
    style_cell(tbl.cell(0, i), h, font_size=13, bold=True, color=NVIDIA_GREEN, bg_color=TABLE_HEADER_BG)

rows = [
    ["Internet", "Required", "Not needed"],
    ["Latency", "100\u2013500ms", "<10ms"],
    ["Disaster zone", "Unusable", "Designed for it"],
    ["Physics validation", "None", "Every action"],
    ["Single point of failure", "Cloud outage", "Self-contained"],
]
for r, row in enumerate(rows):
    bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
    style_cell(tbl.cell(r+1, 0), row[0], font_size=12, color=WHITE, bold=True, bg_color=bg)
    style_cell(tbl.cell(r+1, 1), row[1], font_size=12, color=ACCENT_RED, bg_color=bg)
    style_cell(tbl.cell(r+1, 2), row[2], font_size=12, color=NVIDIA_GREEN, bold=True, bg_color=bg)

tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.0)

# Hardware list (right side)
add_textbox(slide, 7.3, 1.8, 5.5, 0.4, "HARDWARE ON THE TABLE", font_size=14, color=NVIDIA_GREEN, bold=True)

hw_items = [
    ("NVIDIA DGX Spark", "Local LLM \u2014 planning + orchestration"),
    ("NVIDIA Jetson Orin Nano Super", "VLM perception + motor control"),
    ("RTX 3070 Laptop", "Isaac Sim \u2014 physics validation"),
    ("SO-101 Follower Arm", "Executes validated plans"),
    ("SO-101 Leader Arm", "Human override for edge cases"),
    ("2x Mono Cameras", "Stereo perception \u2192 digital twin"),
]

hy = 2.3
for device, role in hw_items:
    add_textbox(slide, 7.3, hy, 5.5, 0.3, device, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, 7.3, hy + 0.3, 5.5, 0.3, role, font_size=12, color=LIGHT_GRAY)
    hy += 0.7

# ============================================================
# SLIDE 8: The Market
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "THE MARKET", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)
add_textbox(slide, 0.8, 1.0, 11.7, 0.7,
    "At the intersection of three rapidly growing markets.",
    font_size=28, color=WHITE, bold=True)

col_data = [
    ("Disaster Response\nRobotics", "$2.5B \u2192 $6.2B", "by 2033 (10.5% CAGR)", "Search & rescue robots:\n$35.3B in 2025\n$70.3B by 2030", ACCENT_RED),
    ("Digital Twin\nTechnology", "$35.8B \u2192 $328.5B", "by 2034 (31.1% CAGR)", "NVIDIA Omniverse deployed\nby TSMC, BMW, Siemens\nfor factory digital twins", ACCENT_BLUE),
    ("Edge AI\nfor Robotics", "$1.2 Trillion", "in 2025 US manufacturing\ninvestment", "NVIDIA Jetson, Isaac Sim,\nDGX Spark all built for\nexactly this use case", NVIDIA_GREEN),
]

x = 0.8
for title, big_num, subtitle, details, color in col_data:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.7), Inches(4.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    add_textbox(slide, x + 0.2, 2.15, 3.3, 0.8, title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.1, 3.3, 0.6, big_num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.8, 3.3, 0.6, subtitle, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 4.6, 3.3, 1.5, details, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += 4.1

add_textbox(slide, 0.8, 6.9, 11.7, 0.4,
    "Sources: Verified Market Reports, Mordor Intelligence, GM Insights, NVIDIA Newsroom",
    font_size=11, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: Beyond Disaster Response
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "BEYOND DISASTER RESPONSE", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)
add_textbox(slide, 0.8, 1.0, 11.7, 0.7,
    "The demo is blocks. The architecture is universal.",
    font_size=28, color=WHITE, bold=True)

applications = [
    ("Hazmat Decommissioning", "Test a cutting procedure in sim before committing a robot in a radioactive environment"),
    ("Remote Space Construction", "Build a habitat in sim on Earth, execute on Mars where 20-min signal delay makes teleoperation impossible"),
    ("Precision Manufacturing", "Simulate micro-assembly before a robot places a component worth thousands of dollars"),
    ("Infrastructure Inspection", "Simulate load redistribution before a robot removes a damaged bridge section"),
]

y = 2.2
for title, desc in applications:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.7), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x50)
    shape.line.width = Pt(1)

    add_textbox(slide, 1.0, y + 0.05, 3.5, 0.5, title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, 1.0, y + 0.5, 10.3, 0.45, desc, font_size=14, color=LIGHT_GRAY)
    y += 1.15

tf = add_textbox(slide, 0.8, 6.5, 11.7, 0.8,
    "The prompt changes. The manipulator changes. The physics engine stays the same.\nIf you can simulate it, you can validate it. If you can validate it, you can trust a robot to do it.",
    font_size=18, color=NVIDIA_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: Live Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)

add_textbox(slide, 1, 2.2, 11.333, 1.2, "LIVE DEMO", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 4.5, 3.5, 4.333)
add_textbox(slide, 1, 3.8, 11.333, 1.0,
    "Build a tower. Break it in simulation.\nRedesign it. Rebuild it stronger.",
    font_size=24, color=NVIDIA_GREEN, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 1, 1.0, 11.333, 1.0, "Edge Rescue", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 4, 2.1, 5.333)

add_textbox(slide, 1.5, 2.5, 10.333, 2.0,
    "An AI system that plans a construction task, simulates it to check for structural failure, and only then tells the robot to build.\n\nIf the simulation says it'll collapse in an earthquake, the AI redesigns it and the robot rebuilds.",
    font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 4.8, 11.333, 0.6,
    "Fully local.  Fully autonomous.  Physics-validated safety.",
    font_size=24, color=NVIDIA_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.8, 11.333, 0.6,
    "Built for the places where cloud doesn't reach and mistakes cost lives.",
    font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.5, 11.333, 0.5,
    "Eashan Iyer  |  Samuel Lihn  |  Gordon Jin  |  Trevor Kwan",
    font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Prize Track Alignment (Appendix)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, "APPENDIX: PRIZE TRACK ALIGNMENT", font_size=14, color=NVIDIA_GREEN, bold=True)
add_accent_line(slide, 0.8, 0.8, 2)

tbl = add_table(slide, 7, 2, 0.8, 1.3, 11.7, 5.5)
style_cell(tbl.cell(0, 0), "Prize", font_size=15, bold=True, color=NVIDIA_GREEN, bg_color=TABLE_HEADER_BG)
style_cell(tbl.cell(0, 1), "Why We Fit", font_size=15, bold=True, color=NVIDIA_GREEN, bg_color=TABLE_HEADER_BG)

prizes = [
    ["NVIDIA Edge AI Track", "DGX Spark + Jetson + Isaac Sim \u2014 the full NVIDIA edge stack in one project"],
    ["NVIDIA Open Models Challenge", "Multi-agent system built on NVIDIA open models and hardware"],
    ["Best Hardware Hack", "3 NVIDIA compute devices + 2 robot arms + stereo cameras on the table"],
    ["Most Technically Complex", "Local LLM + VLA + physics sim + multi-agent orchestration across 3 devices"],
    ["Greylock Best Multi-Turn Agent", "4 agents reasoning about feedback to dynamically complete multi-step construction"],
    ["Grand Prize", "Innovation + technical complexity + social impact demonstrated in a single live demo"],
]
for r, row in enumerate(prizes):
    bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
    style_cell(tbl.cell(r+1, 0), row[0], font_size=14, color=WHITE, bold=True, bg_color=bg)
    style_cell(tbl.cell(r+1, 1), row[1], font_size=14, color=LIGHT_GRAY, bg_color=bg)

tbl.columns[0].width = Inches(4.0)
tbl.columns[1].width = Inches(7.7)

# ============================================================
# SAVE
# ============================================================
output_path = r"C:\Users\Dell\Desktop\treehacks2026\EdgeRescue_TreeHacks2026.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
